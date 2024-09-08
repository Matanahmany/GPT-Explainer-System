"""
presentation_parser.py

This module contains functions to parse PowerPoint files and extract text from slides.
"""

import pptx
import re
from typing import List


def parse(file_path: str) -> List[str]:
    """
    Parses a PowerPoint file and extracts text from slides.

    Args:
        file_path (str): Path to the PowerPoint file.

    Returns:
        List[str]: A list of strings, each containing the text of a slide.
    """
    presentation = pptx.Presentation(file_path)
    slides_text = []

    for slide in presentation.slides:
        slide_text = extract_text_from_slide(slide)
        if slide_text:
            slides_text.append(slide_text)

    return slides_text


def extract_text_from_slide(slide) -> str:
    """
    Extracts text from a single slide.

    Args:
        slide: A slide object from the pptx library.

    Returns:
        str: The combined text of all shapes and tables in the slide.
    """
    slide_text = []

    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = clean_text(shape.text)
            slide_text.append(text)
        elif shape.has_table:
            slide_text.extend(extract_text_from_table(shape.table))

    return "\n".join(slide_text).strip()


def extract_text_from_table(table) -> List[str]:
    """
    Extracts text from a table shape.

    Args:
        table: A table object from the pptx library.

    Returns:
        List[str]: A list of strings, each containing the text of a table cell.
    """
    table_text = []

    for row in table.rows:
        for cell in row.cells:
            text = clean_text(cell.text)
            table_text.append(text)

    return table_text


def clean_text(text: str) -> str:
    """
    Cleans the text by removing extra whitespace.

    Args:
        text (str): The text to be cleaned.

    Returns:
        str: The cleaned text.
    """
    return re.sub(r'\s+', ' ', text.strip())
